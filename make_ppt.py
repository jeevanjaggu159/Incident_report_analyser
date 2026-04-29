import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    # Slide 1: Tools Worked On (Title and Body layout)
    slide_layout = prs.slide_layouts[1]
    slide1 = prs.slides.add_slide(slide_layout)
    title1 = slide1.shapes.title
    body1 = slide1.shapes.placeholders[1]
    
    title1.text = "Tools Worked On"
    
    tf = body1.text_frame
    tf.text = "The following core technologies were used to solve the assignment tasks:"
    
    tools = [
        ("Python & FastAPI", "Built the robust backend REST API, orchestrating the RAG pipeline with high performance."),
        ("Google Gemini AI", "Provided core LLM capabilities to extract structured insight from text and generate semantic embeddings."),
        ("FAISS", "Powered blazing-fast vector similarity searches to retrieve relevant historical context for AI prompts."),
        ("SQLite & SQLAlchemy", "Managed persistence for historical incident data and user accounts locally."),
        ("React.js & Recharts", "Created the interactive, modern user interface and dynamic analytics dashboards."),
        ("Tailwind CSS", "Handled UI styling for a premium, responsive design without custom CSS bloat.")
    ]
    
    for tool, desc in tools:
        p = tf.add_paragraph()
        p.text = f"• {tool}: "
        p.font.bold = True
        p.font.size = Pt(16)
        
        # Add the description part
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(16)

    # Slide 2: Workable Solutions & Architecture
    slide_layout2 = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(slide_layout2)
    title2 = slide2.shapes.title
    body2 = slide2.shapes.placeholders[1]
    
    title2.text = "Workable Solution & Architecture"
    
    tf2 = body2.text_frame
    
    p = tf2.add_paragraph()
    p.text = "Problem Statement: "
    p.font.bold = True
    p.font.size = Pt(16)
    run = p.add_run()
    run.text = "Traffic police needed a streamlined way to instantly process subjective textual incident reports into structured data (Categories, Severities) to identify safety trends like 'No Helmet' violations."
    run.font.bold = False
    
    p = tf2.add_paragraph()
    p.text = "Solution Approach: "
    p.font.bold = True
    p.font.size = Pt(16)
    run = p.add_run()
    run.text = "Developed a Retrieval-Augmented Generation (RAG) pipeline. The system converts reports to vectors (Gemini), searches identical historical incidents (FAISS), and prompts the LLM with this context to reliably classify the report."
    run.font.bold = False

    p = tf2.add_paragraph()
    p.text = "Live App Output: "
    p.font.bold = True
    p.font.size = Pt(16)
    run = p.add_run()
    run.text = "A fully functional web dashboard where users input textual reports, receive instant AI categorizations mapping to UI badges, and visualize time-series crash trends on the Analytics page."
    run.font.bold = False
    
    p = tf2.add_paragraph()
    p.text = "System Architecture: "
    p.font.bold = True
    p.font.size = Pt(16)
    run = p.add_run()
    run.text = "React Frontend ↔ FastAPI Backend ↔ [Gemini API (LLM/Embeddings) + FAISS (Vectors) + DB]"
    run.font.bold = False

    prs.save("Presentation.pptx")
    print("Presentation.pptx generated successfully!")

if __name__ == "__main__":
    create_presentation()
